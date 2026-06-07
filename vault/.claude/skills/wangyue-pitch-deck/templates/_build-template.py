"""Build a clean Wangyue brand PPT template from scratch.

Encodes all extracted brand primitives:
- 16:9 canvas (13.333 x 7.5 in)
- Gold #FFD800 title-underline signature (5.02 x 0.21 in at 0.72, 1.67)
- Deep blue #0045E0 label chips
- Black/gold dots for timeline
- ONE / UNITY / RISK EVERYTHING tagline every content slide
- Title convention "[节] - [副标]"
- 9 slide type templates with placeholders
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

W, H = Inches(13.333), Inches(7.5)
GOLD = RGBColor(0xFF, 0xD8, 0x00)
BLUE = RGBColor(0x00, 0x45, 0xE0)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x0A, 0x0A, 0x14)
GRAY_PANEL = RGBColor(0x20, 0x20, 0x30)
GRAY_LABEL = RGBColor(0x60, 0x60, 0x80)

TU_LEFT, TU_TOP = Inches(0.72), Inches(1.67)
TU_W, TU_H = Inches(5.02), Inches(0.21)
TG_TEXT = "ONE / UNITY / RISK EVERYTHING"
CN_FONT = "玄度"
CN_HEAD = "方正VDL LOGO黑 简 SemiBold"
LATIN = "Barlow SemiBold"

p = Presentation()
p.slide_width, p.slide_height = W, H


def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_dark_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    return bg


def add_gold_underline(slide, left=TU_LEFT, top=TU_TOP, width=TU_W):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, TU_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    return bar


def add_text(slide, text, left, top, width, height, font=CN_FONT, size=16,
             color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_tagline(slide):
    return add_text(slide, TG_TEXT, Inches(8.6), Inches(0.30), Inches(4.5), Inches(0.3),
                    font=LATIN, size=11, color=WHITE, align=PP_ALIGN.RIGHT)


def add_title(slide, text):
    return add_text(slide, text, TU_LEFT, Inches(0.85), Inches(11.0), Inches(0.7),
                    font=CN_HEAD, size=32, color=WHITE, bold=True)


# 1. Cover
s = add_blank(p); add_dark_bg(s)
add_text(s, "[COVER · 主标题]", Inches(0.7), Inches(2.5), Inches(12), Inches(2),
         font=CN_HEAD, size=72, color=GOLD, bold=True)
add_text(s, "[副标 / 项目代号 / 日期]", Inches(0.7), Inches(4.8), Inches(12), Inches(0.6),
         font=CN_FONT, size=22, color=WHITE)
add_text(s, TG_TEXT, Inches(0.7), Inches(6.8), Inches(12), Inches(0.4),
         font=LATIN, size=14, color=GOLD)

# 2. Data hero
s = add_blank(p); add_dark_bg(s); add_tagline(s); add_title(s, "数据 - [SECTION]")
add_gold_underline(s)
add_text(s, "[600 万]", Inches(0.72), Inches(2.3), Inches(6), Inches(1.5),
         font=CN_HEAD, size=72, color=GOLD, bold=True)
add_text(s, "[副数据描述]", Inches(0.72), Inches(3.9), Inches(6), Inches(0.5),
         font=CN_FONT, size=18, color=WHITE)
for i, y in enumerate([Inches(2.82), Inches(4.43), Inches(5.93)]):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), y, Inches(3.94), Inches(0.21))
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background()
    add_text(s, "[BREAKDOWN " + str(i + 1) + ": 平台 / 数字]", Inches(6.5), y - Inches(0.45),
             Inches(5), Inches(0.4), font=CN_FONT, size=18, color=WHITE)

# 3. Timeline
s = add_blank(p); add_dark_bg(s); add_tagline(s); add_title(s, "里程碑 - [时间范围]")
add_gold_underline(s)
xs = [Inches(1.05), Inches(5.09), Inches(8.46), Inches(11.66)]
for x in xs:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(2.7), Inches(0.15), Inches(0.15))
    dot.fill.solid(); dot.fill.fore_color.rgb = GOLD; dot.line.fill.background()
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.77), Inches(10.5), Inches(0.02))
ln.fill.solid(); ln.fill.fore_color.rgb = GOLD; ln.line.fill.background()
for i, x in enumerate(xs):
    add_text(s, "[时段 " + str(i + 1) + "]", x - Inches(0.7), Inches(3.0), Inches(2), Inches(0.5),
             font=LATIN, size=18, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "[事件 " + str(i + 1) + " 标题]", x - Inches(1.2), Inches(3.6), Inches(3), Inches(0.5),
             font=CN_HEAD, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "[事件 " + str(i + 1) + " 描述 12-30 字]", x - Inches(1.2), Inches(4.3),
             Inches(3), Inches(2), font=CN_FONT, size=14, color=WHITE, align=PP_ALIGN.CENTER)

# 4. Section divider
s = add_blank(p); add_dark_bg(s); add_tagline(s); add_title(s, "[SECTION] - 核心卖点")
add_gold_underline(s)
hero = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.31), Inches(1.98), Inches(5.15), Inches(3.77))
hero.fill.solid(); hero.fill.fore_color.rgb = GOLD; hero.line.fill.background()
add_text(s, "[HERO_TITLE\n例如: 中式写实都市]", Inches(1.31), Inches(3.0), Inches(5.15), Inches(1.5),
         font=CN_HEAD, size=44, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "[副标 / 子主题]", Inches(7.0), Inches(3.5), Inches(5.5), Inches(0.6),
         font=CN_FONT, size=24, color=WHITE)

# 5. Content image-heavy
s = add_blank(p); add_dark_bg(s); add_tagline(s); add_title(s, "[节] - [副标]")
add_gold_underline(s)
img_panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(2.3), Inches(7.5), Inches(4.8))
img_panel.fill.solid(); img_panel.fill.fore_color.rgb = GRAY_PANEL; img_panel.line.fill.background()
add_text(s, "[REPLACE: 主图]", Inches(0.72), Inches(4.4), Inches(7.5), Inches(0.5),
         font=CN_FONT, size=18, color=GRAY_LABEL, align=PP_ALIGN.CENTER)
for i, y in enumerate([Inches(2.3), Inches(3.7), Inches(5.1)]):
    chip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), y, Inches(1.59), Inches(0.39))
    chip.fill.solid(); chip.fill.fore_color.rgb = BLUE; chip.line.fill.background()
    add_text(s, "[标签 " + str(i + 1) + "]", Inches(8.5), y + Inches(0.04), Inches(1.59), Inches(0.32),
             font=CN_FONT, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "[Caption " + str(i + 1) + " 注解 1-2 行]", Inches(10.3), y, Inches(2.7), Inches(1.0),
             font=CN_FONT, size=14, color=WHITE)

# 6. Comparison
s = add_blank(p); add_dark_bg(s); add_tagline(s); add_title(s, "[节] - [对比主题]")
add_gold_underline(s)
for i, x in enumerate([Inches(0.72), Inches(6.92)]):
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.3), Inches(5.7), Inches(4.5))
    panel.fill.solid(); panel.fill.fore_color.rgb = GRAY_PANEL; panel.line.fill.background()
    add_text(s, "[LABEL " + str(i + 1) + "]", x, Inches(2.3), Inches(5.7), Inches(0.5),
             font=CN_HEAD, size=22, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "[REPLACE: 图 " + str(i + 1) + "]", x, Inches(4.0), Inches(5.7), Inches(0.6),
             font=CN_FONT, size=16, color=GRAY_LABEL, align=PP_ALIGN.CENTER)
    add_text(s, "[描述 " + str(i + 1) + "]", x, Inches(6.0), Inches(5.7), Inches(0.6),
             font=CN_FONT, size=14, color=WHITE, align=PP_ALIGN.CENTER)

# 7. Feature grid (3-cell)
s = add_blank(p); add_dark_bg(s); add_tagline(s); add_title(s, "[节] - 多 feature 并列")
add_gold_underline(s)
xs = [Inches(0.68), Inches(4.92), Inches(9.15)]
for i, x in enumerate(xs):
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.94), Inches(3.51), Inches(4.04))
    panel.fill.solid(); panel.fill.fore_color.rgb = GRAY_PANEL; panel.line.fill.background()
    add_text(s, "[REPLACE: 图 " + str(i + 1) + "]", x, Inches(4.4), Inches(3.51), Inches(0.5),
             font=CN_FONT, size=14, color=GRAY_LABEL, align=PP_ALIGN.CENTER)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.41), Inches(6.85), Inches(2.68), Inches(0.14))
    line.fill.solid(); line.fill.fore_color.rgb = GOLD; line.line.fill.background()
    add_text(s, "[Feature " + str(i + 1) + " 名]", x, Inches(7.05), Inches(3.51), Inches(0.4),
             font=CN_HEAD, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 8. Event / distribution node
s = add_blank(p); add_dark_bg(s); add_tagline(s); add_title(s, "发行节点 - [事件名]")
add_gold_underline(s, left=Inches(0.13), top=Inches(1.63))
add_text(s, "[事件简介 / 场地 / 时间]", Inches(0.13), Inches(2.3), Inches(6), Inches(0.6),
         font=CN_FONT, size=18, color=WHITE)
add_text(s, "首次 [动作] —— 整体花费超 [X]W", Inches(0.13), Inches(3.0), Inches(6), Inches(0.6),
         font=CN_HEAD, size=22, color=GOLD, bold=True)
banner = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.13), Inches(5.93), Inches(5.74), Inches(0.28))
banner.fill.solid(); banner.fill.fore_color.rgb = GOLD; banner.line.fill.background()
add_text(s, "[嘉宾阵容 / 联动合作]", Inches(0.2), Inches(5.92), Inches(5.6), Inches(0.3),
         font=CN_HEAD, size=14, color=BLACK, bold=True)
for i, y in enumerate([Inches(4.37), Inches(6.82)]):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.17), y, Inches(3.28), Inches(0.32))
    b.fill.solid(); b.fill.fore_color.rgb = GOLD; b.line.fill.background()
    add_text(s, "[联动 / 嘉宾 " + str(i + 1) + "]", Inches(6.25), y, Inches(3.2), Inches(0.32),
             font=CN_HEAD, size=14, color=BLACK, bold=True)
add_text(s, "[REPLACE: 现场图 / coser / KV]", Inches(7.5), Inches(2.3), Inches(5.7), Inches(2),
         font=CN_FONT, size=18, color=GRAY_LABEL, align=PP_ALIGN.CENTER)

# 9. Closing
s = add_blank(p); add_dark_bg(s)
add_text(s, "[致谢 / 联系方式 / 项目联系人]", Inches(0.7), Inches(2.8), Inches(12), Inches(2),
         font=CN_HEAD, size=54, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
add_text(s, TG_TEXT, Inches(0.7), Inches(6.4), Inches(12), Inches(0.6),
         font=LATIN, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

out = r"D:\Administrator\Documents\{{USER_NAME}}\.claude\skills\wangyue-pitch-deck\templates\wangyue-master.pptx"
os.makedirs(os.path.dirname(out), exist_ok=True)
p.save(out)
print("Template saved:", out)
print("Slides:", len(p.slides))
print("Size:", os.path.getsize(out) // 1024, "KB")
